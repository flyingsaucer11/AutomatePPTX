
from pptx import Presentation
from pptx.util import Pt,Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

font_size=13

def generatePPTX(firstName, lastName, idType, idNumber, expiration, dateIssued, dateExpires,img_path,extnum):

    prs = Presentation('production_data/targetPresentationFile.pptx')
    my_slide=prs.slides[0]

    for curr_shape in my_slide.shapes:
        if curr_shape.shape_id==15:
            table=curr_shape.table

            p=table.cell(0,1).text_frame.paragraphs[0]
            myrun=p.add_run()
            myrun.text=idType
            myfont=myrun.font
            myfont.color.rgb=RGBColor(0x00,0x00,0x00)
            myfont.size=Pt(font_size)

            p = table.cell(0, 4).text_frame.paragraphs[0]
            myrun = p.add_run()
            myrun.text = idNumber
            myfont = myrun.font
            myfont.color.rgb = RGBColor(0x00, 0x00, 0x00)
            myfont.size = Pt(font_size)

            p = table.cell(0, 6).text_frame.paragraphs[0]
            myrun = p.add_run()
            myrun.text = expiration
            myfont = myrun.font
            myfont.color.rgb = RGBColor(0x00, 0x00, 0x00)
            myfont.size = Pt(font_size)


        if curr_shape.shape_id==4:
            table=curr_shape.table

            p = table.cell(0, 1).text_frame.paragraphs[0]
            myrun = p.add_run()
            myrun.text = dateIssued
            myfont = myrun.font
            myfont.color.rgb = RGBColor(0x00, 0x00, 0x00)
            myfont.size = Pt(font_size)

            p = table.cell(0, 4).text_frame.paragraphs[0]
            myrun = p.add_run()
            myrun.text = dateIssued
            myfont = myrun.font
            myfont.color.rgb = RGBColor(0x00, 0x00, 0x00)
            myfont.size = Pt(font_size)

        if curr_shape.shape_id==11:
            table=curr_shape.table
            p=table.cell(0,2).text_frame.paragraphs[1]
            p.clear()
            myrun=p.add_run()
            myrun.text=dateExpires
            myfont=myrun.font
            myfont.color.rgb=RGBColor(0x00,0x00,0x00)
            myfont.size = Pt(font_size)

    if extnum==2:
        left=Inches(6.40)
        top=Inches(2.68)
        picheight=Inches(0.9)
    else:
        left=Inches(6.34)
        top=Inches(3.12)
        picheight=Inches(0.4)

    pic=prs.slides[0].shapes.add_picture(img_path,left,top,height=picheight)

    left=Inches(1.25)
    top=Inches(2.80 )
    width=Inches(3.0)
    height=Inches(1.0)
    txBox=my_slide.shapes.add_textbox(left,top,width,height)

    tf=txBox.text_frame
    p=tf.add_paragraph()
    p.text=firstName
    p.font.bold=True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.font.size = Pt(font_size)
    #print('done')
    newfilename=str(idNumber)+'.pptx'
    prs.save(newfilename)

